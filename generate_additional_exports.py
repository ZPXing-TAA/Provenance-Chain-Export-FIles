#!/usr/bin/env python3
"""Generate additional controlled non-AI digital-export image samples.

This CLI covers office/PDF pilot digital-to-digital export classes:

* PDF page rasterization
* DOCX/ODT document export with a local renderer
* XLSX/ODS spreadsheet export with a local renderer
* PPTX/ODP presentation export with a local renderer

It does not call any LLM, image generation model, online API, or remote asset.
Every generated PNG master is a controlled digital-to-digital file export.
"""

from __future__ import annotations

import argparse
import importlib.metadata
import json
import random
import re
import shutil
import subprocess
import sys
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any

from generate_non_ai_exports import (
    CANVAS_CHOICES,
    LOCAL_SAFE_FONTS,
    VISUAL_STYLES,
    WORDLIST_PATH,
    Canvas,
    choose,
    decoded_pixel_hash,
    ensure_empty_or_create,
    number_choice,
    phrase,
    sentence,
    sha256_file,
    sha256_text,
)


SCRIPT_DIR = Path(__file__).resolve().parent

RUN_ORDER = [
    "pdf",
    "document_export",
    "spreadsheet_export",
    "presentation_export",
]

FIXED_DOC_TIME = datetime(2020, 1, 1, 0, 0, 0)
FIXED_ZIP_DT = (2020, 1, 1, 0, 0, 0)

SOURCE_DIRS = {
    "pdf": "pdf",
    "document_export": "documents",
    "spreadsheet_export": "spreadsheets",
    "presentation_export": "presentations",
}

SOURCE_TYPES = {
    "pdf": "pdf",
    "document_export": "docx",
    "spreadsheet_export": "xlsx",
    "presentation_export": "pptx",
}

FAMILIES = {
    "pdf": [
        "pdf_report_page",
        "pdf_notice_page",
        "pdf_form_page",
        "pdf_abstract_page",
    ],
    "document_export": [
        "docx_policy_brief",
        "docx_meeting_notes",
        "docx_research_memo",
        "docx_project_update",
    ],
    "spreadsheet_export": [
        "xlsx_budget_grid",
        "xlsx_inventory_tracker",
        "xlsx_schedule_matrix",
        "xlsx_metric_table",
    ],
    "presentation_export": [
        "pptx_title_content",
        "pptx_two_column",
        "pptx_process_overview",
        "pptx_metric_summary",
    ],
}


def parse_canvas(value: str) -> Canvas:
    match = re.fullmatch(r"(\d+)x(\d+)", value.strip().lower())
    if not match:
        raise argparse.ArgumentTypeError("canvas size must use WIDTHxHEIGHT, e.g. 1024x1024")
    width, height = int(match.group(1)), int(match.group(2))
    if (width, height) not in CANVAS_CHOICES:
        choices = ", ".join(f"{w}x{h}" for w, h in CANVAS_CHOICES)
        raise argparse.ArgumentTypeError(f"canvas must be one of: {choices}")
    return Canvas(width=width, height=height)


def load_wordlists() -> dict[str, Any]:
    with WORDLIST_PATH.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def resolve_source_format(pipeline: str, source_format: str | None) -> str | None:
    allowed = {
        "document_export": ("docx", "odt"),
        "spreadsheet_export": ("xlsx", "ods"),
        "presentation_export": ("pptx", "odp"),
    }
    if pipeline not in allowed:
        if source_format is not None:
            raise ValueError(f"--source-format is not valid for pipeline {pipeline}")
        return None
    if source_format is None:
        return allowed[pipeline][0]
    if source_format not in allowed[pipeline]:
        choices = ", ".join(allowed[pipeline])
        raise ValueError(f"--source-format for {pipeline} must be one of: {choices}")
    return source_format


def relative_string(path: Path | None, base: Path) -> str | None:
    if path is None:
        return None
    try:
        return str(path.resolve().relative_to(base.resolve()))
    except ValueError:
        return str(path)


def package_version(name: str) -> str:
    try:
        return importlib.metadata.version(name)
    except importlib.metadata.PackageNotFoundError:
        return "unknown"


def canvas_for_sample(rng: random.Random, fixed_canvas: Canvas | None) -> Canvas:
    if fixed_canvas:
        return fixed_canvas
    width, height = choose(rng, CANVAS_CHOICES)
    return Canvas(width=width, height=height)


def make_code(rng: random.Random, words: dict[str, Any]) -> str:
    prefix = "".join(rng.choice("ABCDEFGHJKLMNPQRSTUVWXYZ") for _ in range(3))
    return f"{prefix}-{number_choice(rng, words, 'code_blocks')}-{number_choice(rng, words, 'code_pairs')}"


def font_for_sample(rng: random.Random) -> dict[str, str]:
    return choose(rng, LOCAL_SAFE_FONTS)


def make_layout_id(
    pipeline: str,
    family: str,
    style: dict[str, Any],
    canvas: Canvas,
    counts: dict[str, int],
) -> str:
    count_part = "_".join(f"{key}{value}" for key, value in sorted(counts.items()))
    return f"{pipeline}__{family}__{style['style_id']}__{canvas.width}x{canvas.height}__{count_part}"


def record_base(
    pipeline: str,
    index: int,
    base_seed: int,
    rng: random.Random,
    words: dict[str, Any],
    fixed_canvas: Canvas | None,
) -> dict[str, Any]:
    sample_id = f"sample_{index + 1:06d}"
    canvas = canvas_for_sample(rng, fixed_canvas)
    family = FAMILIES[pipeline][index % len(FAMILIES[pipeline])]
    style = VISUAL_STYLES[(index + rng.randrange(len(VISUAL_STYLES))) % len(VISUAL_STYLES)]
    font = font_for_sample(rng)
    return {
        "sample_id": sample_id,
        "pipeline": pipeline,
        "source_type": SOURCE_TYPES[pipeline],
        "template_family": family,
        "style_id": style["style_id"],
        "random_seed": base_seed + index,
        "canvas": canvas,
        "canvas_size": {"width": canvas.width, "height": canvas.height},
        "style": style,
        "font_family_id": font["id"],
        "font_family_css": font["css"],
        "title": phrase(rng, words, 3, 5),
        "subtitle": sentence(rng, words, 8, 12),
        "code": make_code(rng, words),
        "source_path": None,
        "source_asset_hash": None,
        "parent_asset_hashes": None,
        "render_path": None,
        "intermediate_render_path": None,
        "output_path": None,
        "output_asset_hash": None,
        "decoded_pixel_hash": None,
        "renderer": None,
        "renderer_version": None,
    }


def finalize_source(record: dict[str, Any], source_path: Path) -> None:
    record["source_path"] = source_path
    record["source_asset_hash"] = sha256_file(source_path)
    record["parent_asset_hashes"] = [record["source_asset_hash"]]


def normalize_zip_archive(path: Path) -> None:
    """Rewrite a ZIP container with fixed metadata for byte reproducibility."""
    temp_path = path.with_suffix(path.suffix + ".tmp")
    with zipfile.ZipFile(path, "r") as source_zip, zipfile.ZipFile(
        temp_path,
        "w",
        compression=zipfile.ZIP_DEFLATED,
        compresslevel=9,
    ) as target_zip:
        for name in sorted(source_zip.namelist()):
            source_info = source_zip.getinfo(name)
            target_info = zipfile.ZipInfo(filename=name, date_time=FIXED_ZIP_DT)
            target_info.compress_type = zipfile.ZIP_DEFLATED
            target_info.external_attr = source_info.external_attr
            target_info.create_system = source_info.create_system
            payload = source_zip.read(name)
            if name == "docProps/core.xml":
                text = payload.decode("utf-8")
                text = re.sub(
                    r"(<dcterms:created\b[^>]*>).*?(</dcterms:created>)",
                    r"\g<1>2020-01-01T00:00:00Z\2",
                    text,
                )
                text = re.sub(
                    r"(<dcterms:modified\b[^>]*>).*?(</dcterms:modified>)",
                    r"\g<1>2020-01-01T00:00:00Z\2",
                    text,
                )
                payload = text.encode("utf-8")
            elif name == "meta.xml":
                text = payload.decode("utf-8")
                text = re.sub(
                    r"(<meta:creation-date>).*?(</meta:creation-date>)",
                    r"\g<1>2020-01-01T00:00:00Z\2",
                    text,
                )
                text = re.sub(
                    r"(<dc:date>).*?(</dc:date>)",
                    r"\g<1>2020-01-01T00:00:00Z\2",
                    text,
                )
                payload = text.encode("utf-8")
            elif name == "settings.xml":
                text = payload.decode("utf-8")
                text = re.sub(
                    r'(<config:config-item config:name="RsidRoot" config:type="int">).*?(</config:config-item>)',
                    r"\g<1>0\2",
                    text,
                )
                text = re.sub(
                    r'(<config:config-item config:name="Rsid" config:type="int">).*?(</config:config-item>)',
                    r"\g<1>0\2",
                    text,
                )
                payload = text.encode("utf-8")
            target_zip.writestr(target_info, payload)
    temp_path.replace(path)


def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    color = hex_color.strip().lstrip("#")
    return tuple(int(color[idx : idx + 2], 16) for idx in (0, 2, 4))


def make_rows(rng: random.Random, words: dict[str, Any], count: int) -> list[list[str]]:
    rows = []
    for _ in range(count):
        rows.append(
            [
                choose(rng, words["products"]).title(),
                choose(rng, words["departments"]),
                choose(rng, words["status"]),
                str(number_choice(rng, words, "scores")),
                f"{choose(rng, words['months'])[:3]} {number_choice(rng, words, 'days')}",
            ]
        )
    return rows


def write_pdf_source(record: dict[str, Any], rng: random.Random, words: dict[str, Any], output_dir: Path) -> None:
    try:
        from reportlab.lib.colors import HexColor
        from reportlab.pdfgen import canvas as pdf_canvas
    except ModuleNotFoundError as exc:
        raise RuntimeError("reportlab is required for the PDF pipeline") from exc

    source_dir = output_dir / SOURCE_DIRS["pdf"]
    source_dir.mkdir(parents=True, exist_ok=True)
    path = source_dir / f"{record['sample_id']}.pdf"
    canvas_size = record["canvas"]
    style = record["style"]
    family = record["template_family"]
    rows = make_rows(rng, words, rng.randint(6, 10))
    counts = {"rows": len(rows), "sections": rng.randint(2, 4), "charts": 1 if "report" in family else 0}
    record["layout_counts"] = counts
    record["layout_id"] = make_layout_id("pdf", family, style, canvas_size, counts)

    c = pdf_canvas.Canvas(str(path), pagesize=(canvas_size.width, canvas_size.height), invariant=1)
    w, h = canvas_size.width, canvas_size.height
    margin = 52 if min(w, h) > 800 else 34
    c.setFillColor(HexColor(style["bg"]))
    c.rect(0, 0, w, h, fill=1, stroke=0)
    c.setFillColor(HexColor(style["panel"]))
    c.roundRect(margin, margin, w - 2 * margin, h - 2 * margin, 8, fill=1, stroke=0)
    c.setStrokeColor(HexColor(style["line"]))
    c.setLineWidth(1)
    c.line(margin + 24, h - margin - 92, w - margin - 24, h - margin - 92)

    title_font = "Times-Bold" if style["style_id"] in {"academic_paper", "retro_newspaper"} else "Helvetica-Bold"
    body_font = "Times-Roman" if style["style_id"] in {"academic_paper", "retro_newspaper"} else "Helvetica"
    c.setFillColor(HexColor(style["ink"]))
    c.setFont(title_font, 28 if min(w, h) <= 760 else 38)
    c.drawString(margin + 28, h - margin - 54, record["title"][:58])
    c.setFont(body_font, 11 if min(w, h) <= 760 else 14)
    c.setFillColor(HexColor(style["muted"]))
    c.drawString(margin + 30, h - margin - 78, f"{choose(rng, words['organizations'])} / {record['code']}")

    y = h - margin - 130
    c.setFillColor(HexColor(style["ink"]))
    c.setFont(body_font, 12 if min(w, h) <= 760 else 15)
    for line in [record["subtitle"], sentence(rng, words, 12, 18), sentence(rng, words, 10, 16)]:
        c.drawString(margin + 30, y, line[:100])
        y -= 24

    table_y = max(margin + 120, y - 32)
    col_x = [margin + 32, margin + 230, margin + 395, margin + 530, margin + 640]
    if w < 900:
        col_x = [margin + 24, margin + 170, margin + 292, margin + 392, margin + 482]
    c.setFont("Helvetica-Bold", 10 if w < 900 else 12)
    c.setFillColor(HexColor(style["accent"]))
    for idx, heading in enumerate(["Item", "Owner", "Status", "Score", "Date"]):
        c.drawString(col_x[idx], table_y, heading)
    c.setFont(body_font, 9 if w < 900 else 11)
    c.setFillColor(HexColor(style["ink"]))
    row_y = table_y - 22
    for row in rows[:8]:
        c.setStrokeColor(HexColor(style["line"]))
        c.line(margin + 28, row_y - 7, w - margin - 28, row_y - 7)
        for idx, cell in enumerate(row):
            c.drawString(col_x[idx], row_y, cell[:22])
        row_y -= 24

    if counts["charts"]:
        c.setFillColor(HexColor(style["accent2"]))
        chart_x = margin + 32
        chart_y = margin + 54
        bar_w = max(26, (w - 2 * margin - 80) / 8)
        for idx in range(8):
            val = int(number_choice(rng, words, "chart_values"))
            c.rect(chart_x + idx * (bar_w + 8), chart_y, bar_w, val * 1.6, fill=1, stroke=0)
        c.setFillColor(HexColor(style["muted"]))
        c.setFont(body_font, 10)
        c.drawString(chart_x, chart_y - 18, choose(rng, words["metrics"]))

    c.setFillColor(HexColor(style["muted"]))
    c.setFont(body_font, 9)
    c.drawRightString(w - margin - 30, margin + 24, "self-authored digital PDF source")
    c.save()
    finalize_source(record, path)


def rasterize_pdf_records(records: list[dict[str, Any]], output_dir: Path) -> None:
    try:
        import fitz
    except ModuleNotFoundError as exc:
        raise RuntimeError("PyMuPDF is required for PDF rasterization") from exc
    if hasattr(fitz, "TOOLS"):
        fitz.TOOLS.mupdf_display_errors(False)

    image_dir = output_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)
    renderer_version = f"pymupdf {package_version('PyMuPDF')}"
    for record in records:
        output_path = image_dir / f"{record['sample_id']}.png"
        with fitz.open(str(record["source_path"])) as document:
            page = document[0]
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1, 1), alpha=False)
            pixmap.save(str(output_path))
        record["output_path"] = output_path
        record["output_asset_hash"] = sha256_file(output_path)
        record["decoded_pixel_hash"] = decoded_pixel_hash(output_path)
        record["renderer"] = "pymupdf_pdf_rasterizer"
        record["renderer_version"] = renderer_version


def write_docx_source(record: dict[str, Any], rng: random.Random, words: dict[str, Any], output_dir: Path) -> None:
    try:
        from docx import Document
        from docx.shared import Inches, Pt
    except ModuleNotFoundError as exc:
        raise RuntimeError("python-docx is required for the DOCX pipeline") from exc

    source_dir = output_dir / SOURCE_DIRS["document_export"]
    source_dir.mkdir(parents=True, exist_ok=True)
    path = source_dir / f"{record['sample_id']}.docx"
    canvas_size = record["canvas"]
    rows = make_rows(rng, words, rng.randint(4, 7))
    counts = {"paragraphs": rng.randint(4, 7), "rows": len(rows), "tables": 1}
    record["layout_counts"] = counts
    record["layout_id"] = make_layout_id("docx", record["template_family"], record["style"], canvas_size, counts)

    doc = Document()
    doc.core_properties.author = "controlled-local-generator"
    doc.core_properties.created = FIXED_DOC_TIME
    doc.core_properties.modified = FIXED_DOC_TIME
    doc.core_properties.last_modified_by = "controlled-local-generator"
    section = doc.sections[0]
    section.page_width = Inches(min(13.0, max(7.5, canvas_size.width / 96)))
    section.page_height = Inches(min(18.0, max(7.5, canvas_size.height / 96)))
    margin = Inches(0.55)
    section.top_margin = margin
    section.bottom_margin = margin
    section.left_margin = margin
    section.right_margin = margin

    title = doc.add_heading(record["title"], level=0)
    for run in title.runs:
        run.font.name = "Arial"
        run.font.size = Pt(24)
    subtitle = doc.add_paragraph(f"{choose(rng, words['organizations'])} / {record['code']}")
    subtitle.runs[0].font.size = Pt(10)
    doc.add_paragraph(record["subtitle"])
    for _ in range(counts["paragraphs"]):
        doc.add_paragraph(sentence(rng, words, 14, 22))

    table = doc.add_table(rows=1, cols=5)
    table.style = "Table Grid"
    for idx, heading in enumerate(["Item", "Owner", "Status", "Score", "Date"]):
        table.rows[0].cells[idx].text = heading
    for row in rows:
        cells = table.add_row().cells
        for idx, cell in enumerate(row):
            cells[idx].text = cell
    doc.add_paragraph("Local DOCX source generated by deterministic rules.")
    doc.save(path)
    normalize_zip_archive(path)
    finalize_source(record, path)


def libreoffice_version(soffice: str) -> str:
    try:
        result = subprocess.run(
            [soffice, "--version"],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
    except (OSError, subprocess.CalledProcessError):
        return "libreoffice unknown"
    return result.stdout.strip() or result.stderr.strip() or "libreoffice unknown"


def convert_docx_with_libreoffice(records: list[dict[str, Any]], output_dir: Path, soffice: str) -> str:
    return convert_office_sources_with_libreoffice(records, output_dir, soffice)


def convert_office_sources_with_libreoffice(records: list[dict[str, Any]], output_dir: Path, soffice: str) -> str:
    pdf_dir = output_dir / "pdf_intermediate"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    version = libreoffice_version(soffice)
    for record in records:
        subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(pdf_dir),
                str(record["source_path"]),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        pdf_path = pdf_dir / f"{Path(record['source_path']).stem}.pdf"
        if not pdf_path.exists():
            raise RuntimeError(f"LibreOffice did not create expected PDF: {pdf_path}")
        record["render_path"] = pdf_path
        record["intermediate_render_path"] = pdf_path
    return version


def convert_primary_source_with_libreoffice(
    record: dict[str, Any],
    output_dir: Path,
    soffice: str,
    target_format: str,
) -> None:
    source_path = Path(record["source_path"])
    source_dir = output_dir / SOURCE_DIRS[record["pipeline"]]
    source_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            target_format,
            "--outdir",
            str(source_dir),
            str(source_path),
        ],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    converted_path = source_dir / f"{source_path.stem}.{target_format}"
    if not converted_path.exists():
        raise RuntimeError(f"LibreOffice did not create expected {target_format.upper()} source: {converted_path}")
    normalize_zip_archive(converted_path)
    if converted_path != source_path and source_path.exists():
        source_path.unlink()
    finalize_source(record, converted_path)
    record["source_type"] = target_format


def rasterize_office_pdf_records(
    records: list[dict[str, Any]],
    output_dir: Path,
    libreoffice_version_text: str,
    renderer: str,
) -> None:
    try:
        import fitz
    except ModuleNotFoundError as exc:
        raise RuntimeError("PyMuPDF is required for LibreOffice PDF rasterization") from exc
    if hasattr(fitz, "TOOLS"):
        fitz.TOOLS.mupdf_display_errors(False)

    image_dir = output_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)
    renderer_version = f"{libreoffice_version_text}; pymupdf {package_version('PyMuPDF')}"
    for record in records:
        output_path = image_dir / f"{record['sample_id']}.png"
        with fitz.open(str(record["render_path"])) as document:
            page = document[0]
            target = record["canvas_size"]
            scale_x = target["width"] / page.rect.width
            scale_y = target["height"] / page.rect.height
            pixmap = page.get_pixmap(matrix=fitz.Matrix(scale_x, scale_y), alpha=False)
            pixmap.save(str(output_path))
        record["output_path"] = output_path
        record["output_asset_hash"] = sha256_file(output_path)
        record["decoded_pixel_hash"] = decoded_pixel_hash(output_path)
        record["renderer"] = renderer
        record["renderer_version"] = renderer_version


def write_xlsx_source(record: dict[str, Any], rng: random.Random, words: dict[str, Any], output_dir: Path) -> None:
    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
        from openpyxl.utils import get_column_letter
    except ModuleNotFoundError as exc:
        raise RuntimeError("openpyxl is required for the spreadsheet export pipeline") from exc

    source_dir = output_dir / SOURCE_DIRS["spreadsheet_export"]
    source_dir.mkdir(parents=True, exist_ok=True)
    path = source_dir / f"{record['sample_id']}.xlsx"
    rows = make_rows(rng, words, rng.randint(10, 15))
    counts = {"rows": len(rows), "columns": 5, "formulas": 3, "charts": 1}
    record["layout_counts"] = counts
    record["layout_id"] = make_layout_id("spreadsheet_export", record["template_family"], record["style"], record["canvas"], counts)

    style = record["style"]
    wb = Workbook()
    wb.properties.creator = "controlled-local-generator"
    wb.properties.lastModifiedBy = "controlled-local-generator"
    wb.properties.created = FIXED_DOC_TIME
    wb.properties.modified = FIXED_DOC_TIME
    ws = wb.active
    ws.title = "Export"
    ws.sheet_view.showGridLines = False

    headers = ["Item", "Owner", "Status", "Score", "Date"]
    ws["A1"] = record["title"]
    ws["A2"] = f"{choose(rng, words['organizations'])} / {record['code']}"
    ws["A3"] = record["subtitle"]
    ws.merge_cells("A1:E1")
    ws.merge_cells("A2:E2")
    ws.merge_cells("A3:E3")

    ws["A5"] = "Summary"
    ws["B5"] = "=AVERAGE(D8:D99)"
    ws["C5"] = "=MAX(D8:D99)"
    ws["D5"] = "=MIN(D8:D99)"
    ws["E5"] = "XLSX export"
    ws.append([])
    ws.append(headers)
    for row in rows:
        ws.append([row[0], row[1], row[2], int(row[3]), row[4]])

    accent = style["accent"].lstrip("#")
    accent2 = style["accent2"].lstrip("#")
    panel = style["panel"].lstrip("#")
    ink = style["ink"].lstrip("#")
    line = style["line"].lstrip("#")
    muted = style["muted"].lstrip("#")
    thin = Side(style="thin", color=line)

    ws["A1"].font = Font(name="Arial", size=22, bold=True, color=ink)
    ws["A2"].font = Font(name="Arial", size=10, color=muted)
    ws["A3"].font = Font(name="Arial", size=11, color=ink)
    for cell in ws[5]:
        cell.fill = PatternFill("solid", fgColor=panel)
        cell.border = Border(bottom=thin)
        cell.alignment = Alignment(vertical="center")
    for cell in ws[7]:
        cell.fill = PatternFill("solid", fgColor=accent)
        cell.font = Font(name="Arial", bold=True, color="FFFFFF")
        cell.alignment = Alignment(horizontal="left")
    for row in ws.iter_rows(min_row=8, max_row=7 + len(rows), min_col=1, max_col=5):
        for cell in row:
            cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
            cell.alignment = Alignment(vertical="center")
        if row[0].row % 2 == 0:
            for cell in row:
                cell.fill = PatternFill("solid", fgColor=accent2)
    for col, width in zip(range(1, 6), [24, 22, 16, 12, 14]):
        ws.column_dimensions[get_column_letter(col)].width = width
    ws.row_dimensions[1].height = 32
    ws.row_dimensions[3].height = 36
    ws.freeze_panes = "A8"
    ws.auto_filter.ref = f"A7:E{7 + len(rows)}"

    chart = BarChart()
    chart.title = choose(rng, words["metrics"])
    chart.y_axis.title = "Score"
    chart.x_axis.title = "Item"
    data = Reference(ws, min_col=4, min_row=7, max_row=7 + len(rows))
    cats = Reference(ws, min_col=1, min_row=8, max_row=7 + len(rows))
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 6
    chart.width = 12
    ws.add_chart(chart, "G2")

    ws.page_setup.orientation = "landscape" if record["canvas"].width >= record["canvas"].height else "portrait"
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 1
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.print_area = f"A1:L{max(22, 9 + len(rows))}"
    ws.page_margins.left = 0.3
    ws.page_margins.right = 0.3
    ws.page_margins.top = 0.35
    ws.page_margins.bottom = 0.35
    wb.save(path)
    normalize_zip_archive(path)
    finalize_source(record, path)


def write_pptx_source(record: dict[str, Any], rng: random.Random, words: dict[str, Any], output_dir: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ModuleNotFoundError as exc:
        raise RuntimeError("python-pptx is required for the presentation export pipeline") from exc

    source_dir = output_dir / SOURCE_DIRS["presentation_export"]
    source_dir.mkdir(parents=True, exist_ok=True)
    path = source_dir / f"{record['sample_id']}.pptx"
    bullets = [sentence(rng, words, 6, 10) for _ in range(rng.randint(3, 5))]
    counts = {"bullets": len(bullets), "metrics": 3, "shapes": rng.randint(4, 7)}
    record["layout_counts"] = counts
    record["layout_id"] = make_layout_id("presentation_export", record["template_family"], record["style"], record["canvas"], counts)

    style = record["style"]
    prs = Presentation()
    prs.core_properties.author = "controlled-local-generator"
    prs.core_properties.created = FIXED_DOC_TIME
    prs.core_properties.modified = FIXED_DOC_TIME
    prs.core_properties.last_modified_by = "controlled-local-generator"
    prs.slide_width = Inches(record["canvas"].width / 96)
    prs.slide_height = Inches(record["canvas"].height / 96)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def color(hex_color: str) -> RGBColor:
        r, g, b = hex_to_rgb(hex_color)
        return RGBColor(r, g, b)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color(style["bg"])
    sw = prs.slide_width
    sh = prs.slide_height
    margin = int(sw * 0.065)

    for _ in range(counts["shapes"]):
        x = rng.randint(int(sw * 0.08), int(sw * 0.82))
        y = rng.randint(int(sh * 0.12), int(sh * 0.78))
        size = rng.randint(int(sw * 0.03), int(sw * 0.09))
        shape_type = choose(rng, [MSO_SHAPE.RECTANGLE, MSO_SHAPE.OVAL, MSO_SHAPE.DIAMOND])
        shape = slide.shapes.add_shape(shape_type, x, y, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(style["accent2"])
        shape.fill.transparency = 72
        shape.line.color.rgb = color(style["line"])

    eyebrow = slide.shapes.add_textbox(margin, int(sh * 0.05), int(sw * 0.72), int(sh * 0.04))
    eyebrow_tf = eyebrow.text_frame
    eyebrow_tf.text = f"{choose(rng, words['departments']).upper()} / {record['code']}"
    eyebrow_run = eyebrow_tf.paragraphs[0].runs[0]
    eyebrow_run.font.name = "Arial"
    eyebrow_run.font.size = Pt(14)
    eyebrow_run.font.bold = True
    eyebrow_run.font.color.rgb = color(style["accent"])

    title = slide.shapes.add_textbox(margin, int(sh * 0.10), int(sw * 0.76), int(sh * 0.20))
    title_tf = title.text_frame
    title_tf.word_wrap = True
    title_tf.text = record["title"]
    title_run = title_tf.paragraphs[0].runs[0]
    title_run.font.name = "Arial"
    title_run.font.size = Pt(44 if record["canvas"].width >= record["canvas"].height else 34)
    title_run.font.bold = True
    title_run.font.color.rgb = color(style["ink"])

    subtitle = slide.shapes.add_textbox(margin, int(sh * 0.30), int(sw * 0.66), int(sh * 0.10))
    subtitle_tf = subtitle.text_frame
    subtitle_tf.word_wrap = True
    subtitle_tf.text = record["subtitle"]
    subtitle_run = subtitle_tf.paragraphs[0].runs[0]
    subtitle_run.font.name = "Arial"
    subtitle_run.font.size = Pt(18)
    subtitle_run.font.color.rgb = color(style["muted"])

    bullet_box = slide.shapes.add_textbox(margin, int(sh * 0.44), int(sw * 0.58), int(sh * 0.28))
    tf = bullet_box.text_frame
    tf.clear()
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = color(style["ink"])

    card_w = int(sw * 0.26)
    card_h = int(sh * 0.09)
    for idx in range(3):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(sw * 0.66),
            int(sh * (0.50 + idx * 0.12)),
            card_w,
            card_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color(style["panel"])
        card.line.color.rgb = color(style["line"])
        text_frame = card.text_frame
        text_frame.text = f"{number_choice(rng, words, 'percentages')}  {choose(rng, words['metrics'])}"
        run = text_frame.paragraphs[0].runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = color(style["accent"])

    prs.save(path)
    normalize_zip_archive(path)
    finalize_source(record, path)


def build_records(
    pipeline: str,
    count: int,
    base_seed: int,
    fixed_canvas: Canvas | None,
    output_dir: Path,
    words: dict[str, Any],
    source_format: str | None,
) -> list[dict[str, Any]]:
    records = []
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    for index in range(count):
        rng = random.Random(base_seed + index)
        record = record_base(pipeline, index, base_seed, rng, words, fixed_canvas)
        if pipeline == "pdf":
            write_pdf_source(record, rng, words, output_dir)
        elif pipeline == "document_export":
            write_docx_source(record, rng, words, output_dir)
            record["source_type"] = "docx"
            if source_format == "odt":
                if not soffice:
                    raise RuntimeError("LibreOffice/soffice is required to create ODT sources")
                convert_primary_source_with_libreoffice(record, output_dir, soffice, "odt")
        elif pipeline == "spreadsheet_export":
            write_xlsx_source(record, rng, words, output_dir)
            record["source_type"] = "xlsx"
            if source_format == "ods":
                if not soffice:
                    raise RuntimeError("LibreOffice/soffice is required to create ODS sources")
                convert_primary_source_with_libreoffice(record, output_dir, soffice, "ods")
        elif pipeline == "presentation_export":
            write_pptx_source(record, rng, words, output_dir)
            record["source_type"] = "pptx"
            if source_format == "odp":
                if not soffice:
                    raise RuntimeError("LibreOffice/soffice is required to create ODP sources")
                convert_primary_source_with_libreoffice(record, output_dir, soffice, "odp")
        else:
            raise ValueError(f"unknown pipeline: {pipeline}")
        records.append(record)
    return records


def render_records(
    pipeline: str,
    records: list[dict[str, Any]],
    output_dir: Path,
) -> None:
    if pipeline == "pdf":
        rasterize_pdf_records(records, output_dir)
    elif pipeline == "document_export":
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            raise RuntimeError("LibreOffice/soffice renderer requested but not found in PATH")
        libreoffice_version_text = convert_docx_with_libreoffice(records, output_dir, soffice)
        source_type = records[0]["source_type"] if records else "docx"
        rasterize_office_pdf_records(
            records,
            output_dir,
            libreoffice_version_text,
            f"libreoffice_{source_type}_pdf_pymupdf_rasterizer",
        )
    elif pipeline == "spreadsheet_export":
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            raise RuntimeError("LibreOffice/soffice renderer requested but not found in PATH")
        libreoffice_version_text = convert_office_sources_with_libreoffice(records, output_dir, soffice)
        source_type = records[0]["source_type"] if records else "xlsx"
        rasterize_office_pdf_records(
            records,
            output_dir,
            libreoffice_version_text,
            f"libreoffice_calc_{source_type}_pdf_pymupdf_rasterizer",
        )
    elif pipeline == "presentation_export":
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            raise RuntimeError("LibreOffice/soffice renderer requested but not found in PATH")
        libreoffice_version_text = convert_office_sources_with_libreoffice(records, output_dir, soffice)
        source_type = records[0]["source_type"] if records else "pptx"
        rasterize_office_pdf_records(
            records,
            output_dir,
            libreoffice_version_text,
            f"libreoffice_impress_{source_type}_pdf_pymupdf_rasterizer",
        )
    else:
        raise ValueError(f"unknown pipeline: {pipeline}")


def manifest_from_record(record: dict[str, Any], output_dir: Path, wordlist_version: str) -> dict[str, Any]:
    manifest = {
        "schema_version": "provenance_chain_v1",
        "sample_id": record["sample_id"],
        "source_type": record["source_type"],
        "source_path": relative_string(record["source_path"], output_dir),
        "source_asset_hash": record["source_asset_hash"],
        "parent_asset_hashes": record["parent_asset_hashes"],
        "template_family": record["template_family"],
        "style_id": record["style_id"],
        "layout_id": record["layout_id"],
        "random_seed": record["random_seed"],
        "renderer": record["renderer"],
        "renderer_version": record["renderer_version"],
        "canvas_size": record["canvas_size"],
        "output_format": "png",
        "output_path": relative_string(record["output_path"], output_dir),
        "output_asset_hash": record["output_asset_hash"],
        "decoded_pixel_hash": record["decoded_pixel_hash"],
        "font_family_id": record["font_family_id"],
        "layout_counts": record["layout_counts"],
        "wordlist_version": wordlist_version,
        "label_status": "controlled_complete",
        "initial_digital_event": "digital_to_digital_file_export",
        "last_acquisition_event": "none",
        "chain_events": [],
        "ai_inference_used": False,
        "external_assets_used": False,
    }
    if record.get("intermediate_render_path"):
        manifest["intermediate_render_path"] = relative_string(record["intermediate_render_path"], output_dir)
    if record.get("mask_paths"):
        manifest["mask_paths"] = [relative_string(path, output_dir) for path in record["mask_paths"]]
        manifest["mask_asset_hashes"] = record.get("mask_asset_hashes", [])
    if record.get("source_region_ids"):
        manifest["source_region_ids"] = record["source_region_ids"]
    if record.get("spliced_region_mask_path"):
        manifest["spliced_region_mask_path"] = relative_string(record["spliced_region_mask_path"], output_dir)
        manifest["spliced_region_mask_hash"] = record.get("spliced_region_mask_hash")
    return manifest


def write_manifests(records: list[dict[str, Any]], output_dir: Path, wordlist_version: str) -> None:
    manifest_dir = output_dir / "manifests"
    manifest_dir.mkdir(parents=True, exist_ok=True)
    with (manifest_dir / "samples.jsonl").open("w", encoding="utf-8") as handle:
        for record in records:
            manifest = manifest_from_record(record, output_dir, wordlist_version)
            (manifest_dir / f"{record['sample_id']}.json").write_text(
                json.dumps(manifest, indent=2, sort_keys=True) + "\n",
                encoding="utf-8",
            )
            handle.write(json.dumps(manifest, sort_keys=True) + "\n")


def write_summary(pipeline: str, records: list[dict[str, Any]], output_dir: Path) -> None:
    qa_dir = output_dir / "qa"
    qa_dir.mkdir(parents=True, exist_ok=True)
    summary = {
        "pipeline": pipeline,
        "sample_count": len(records),
        "source_type_counts": {},
        "template_family_counts": {},
        "style_counts": {},
        "canvas_counts": {},
        "renderer_counts": {},
        "initial_digital_event": "digital_to_digital_file_export",
        "chain_events": [],
        "ai_inference_used": False,
        "external_assets_used": False,
    }
    for record in records:
        for key, source in [
            ("source_type_counts", record["source_type"]),
            ("template_family_counts", record["template_family"]),
            ("style_counts", record["style_id"]),
            ("canvas_counts", f"{record['canvas_size']['width']}x{record['canvas_size']['height']}"),
            ("renderer_counts", record["renderer"]),
        ]:
            summary[key][source] = summary[key].get(source, 0) + 1
    (qa_dir / "summary.json").write_text(json.dumps(summary, indent=2, sort_keys=True) + "\n", encoding="utf-8")


def run_pipeline(
    pipeline: str,
    count: int,
    seed: int,
    fixed_canvas: Canvas | None,
    output_dir: Path,
    overwrite: bool,
    source_format: str | None = None,
) -> None:
    source_format = resolve_source_format(pipeline, source_format)
    ensure_empty_or_create(output_dir, overwrite)
    for child in {SOURCE_DIRS[pipeline], "images", "manifests", "qa"}:
        (output_dir / child).mkdir(parents=True, exist_ok=True)
    words = load_wordlists()
    records = build_records(pipeline, count, seed, fixed_canvas, output_dir, words, source_format)
    render_records(pipeline, records, output_dir)
    write_manifests(records, output_dir, words.get("version", "unknown"))
    write_summary(pipeline, records, output_dir)
    print(f"{pipeline}: generated {len(records)} PNG masters")
    print(f"{pipeline}: manifest JSONL {output_dir / 'manifests' / 'samples.jsonl'}")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--pipeline", choices=RUN_ORDER + ["all"], required=True)
    parser.add_argument("--count", type=int, default=20, help="Samples per selected pipeline.")
    parser.add_argument("--seed", type=int, default=20260615)
    parser.add_argument("--canvas", type=parse_canvas, default=None, help="Optional fixed canvas, e.g. 1024x1024.")
    parser.add_argument("--output-dir", type=Path, required=True)
    parser.add_argument("--overwrite", action="store_true")
    parser.add_argument(
        "--source-format",
        choices=["docx", "odt", "xlsx", "ods", "pptx", "odp"],
        default=None,
        help="Office source format for document/spreadsheet/presentation exports.",
    )
    args = parser.parse_args(argv)

    if args.count < 1:
        parser.error("--count must be >= 1")
    if args.pipeline == "all" and args.source_format is not None:
        parser.error("--source-format is only valid with a single office pipeline")

    pipelines = RUN_ORDER if args.pipeline == "all" else [args.pipeline]
    for offset, pipeline in enumerate(pipelines):
        pipeline_output = args.output_dir / pipeline if args.pipeline == "all" else args.output_dir
        run_pipeline(
            pipeline=pipeline,
            count=args.count,
            seed=args.seed + offset * 100000,
            fixed_canvas=args.canvas,
            output_dir=pipeline_output,
            overwrite=args.overwrite,
            source_format=args.source_format,
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
